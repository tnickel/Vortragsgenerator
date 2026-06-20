import os
import sys
import json
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Theme Colors (matching references and design specifications)
COLORS = {
    "ink": RGBColor(0x1E, 0x21, 0x52),       # Dark ink for text
    "violet": RGBColor(0x5B, 0x5F, 0xE0),    # Accent purple
    "violetL": RGBColor(0x8E, 0x92, 0xF0),   # Light purple for footers/lines
    "coral": RGBColor(0xFF, 0x7A, 0x6B),     # Warm coral accents
    "cream": RGBColor(0xFF, 0xF8, 0xEE),     # Cream background
    "white": RGBColor(0xFF, 0xFF, 0xFF)
}

def create_presentation(data_json, output_path, images_dir):
    prs = Presentation()
    
    # Set slide dimensions to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Use blank slide layout (index 6 is typically a blank layout)
    blank_layout = prs.slide_layouts[6]
    
    title_text = data_json.get("title", "Präsentation")
    subtitle_text = data_json.get("subtitle", "Erstellt mit Antigravity Vortragsgenerator")
    slides_data = data_json.get("slides", [])
    
    for idx, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        
        # Apply Solid Background Color (Cream)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS["cream"]
        
        # 1. Slide Footer (Title Left, Slide Number Right)
        footer_title_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(6.0), Inches(0.3))
        tf_footer_l = footer_title_box.text_frame
        tf_footer_l.word_wrap = True
        p_footer_l = tf_footer_l.paragraphs[0]
        p_footer_l.text = title_text
        p_footer_l.font.name = "Arial"
        p_footer_l.font.size = Pt(10)
        p_footer_l.font.color.rgb = COLORS["violetL"]
        
        footer_num_box = slide.shapes.add_textbox(Inches(10.5), Inches(6.9), Inches(2.2), Inches(0.3))
        tf_footer_r = footer_num_box.text_frame
        tf_footer_r.word_wrap = True
        p_footer_r = tf_footer_r.paragraphs[0]
        p_footer_r.text = f"Folie {idx + 1} von {len(slides_data)}"
        p_footer_r.alignment = PP_ALIGN.RIGHT
        p_footer_r.font.name = "Arial"
        p_footer_r.font.size = Pt(10)
        p_footer_r.font.color.rgb = COLORS["violetL"]
        
        # 2. Add Speaker notes (for ElevenLabs synthesis)
        notes = slide_data.get("notes", "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
            
        # 3. Slide Content Layout
        is_title_slide = (idx == 0 or slide_data.get("type") == "title")
        
        if is_title_slide:
            # Title slide layout
            
            # Left sidebar decoration (Solid Coral block)
            decor = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(0.4), Inches(7.5))
            decor.fill.solid()
            decor.fill.fore_color.rgb = COLORS["coral"]
            decor.line.fill.background()
            
            # Title & Subtitle Container Box
            text_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(6.2), Inches(3.5))
            tf = text_box.text_frame
            tf.word_wrap = True
            
            # Title paragraph
            p1 = tf.paragraphs[0]
            p1.text = slide_data.get("title", title_text)
            p1.font.name = "Arial"
            p1.font.size = Pt(44)
            p1.font.bold = True
            p1.font.color.rgb = COLORS["ink"]
            p1.space_after = Pt(20)
            
            # Subtitle paragraph
            subtitle = slide_data.get("subtitle", subtitle_text)
            if subtitle:
                p2 = tf.add_paragraph()
                p2.text = subtitle
                p2.font.name = "Arial"
                p2.font.size = Pt(20)
                p2.font.color.rgb = COLORS["violet"]
                
            # Image on the right
            image_name = slide_data.get("image_name", "mascot_hello")
            has_ext = any(image_name.lower().endswith(ext) for ext in [".png", ".jpg", ".jpeg", ".gif"])
            if not has_ext:
                image_name += ".png"
            
            project_dir = os.path.dirname(os.path.abspath(output_path))
            custom_image_path = os.path.join(project_dir, "custom_images", image_name)
            if os.path.exists(custom_image_path):
                image_path = custom_image_path
            else:
                image_path = os.path.join(images_dir, image_name)
            
            if os.path.exists(image_path):
                slide.shapes.add_picture(image_path, Inches(7.2), Inches(1.2), width=Inches(5.0), height=Inches(5.0))
                
        else:
            # Standard slide layout
            
            # Slide Header Title
            title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.0), Inches(0.8))
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            p_title.text = slide_data.get("title", "")
            p_title.font.name = "Arial"
            p_title.font.size = Pt(32)
            p_title.font.bold = True
            p_title.font.color.rgb = COLORS["ink"]
            
            # Top Bar Accent Line (Coral)
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.25), Inches(5.8), Inches(0.04))
            line.fill.solid()
            line.fill.fore_color.rgb = COLORS["coral"]
            line.line.fill.background()
            
            # Bullet Points
            bullets = slide_data.get("bullets", [])
            if isinstance(bullets, str):
                try:
                    bullets = json.loads(bullets)
                except Exception:
                    bullets = [bullets]
                    
            bullets_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(5.8), Inches(4.8))
            tf_bullets = bullets_box.text_frame
            tf_bullets.word_wrap = True
            
            for b_idx, bullet_text in enumerate(bullets):
                p = tf_bullets.paragraphs[0] if b_idx == 0 else tf_bullets.add_paragraph()
                p.text = bullet_text
                p.level = 0
                p.font.name = "Arial"
                p.font.size = Pt(16)
                p.font.color.rgb = COLORS["ink"]
                # Simulating line spacing and paragraph spacing
                p.space_after = Pt(12)
                p.space_before = Pt(6)
                
            # Right side cartoon image
            image_name = slide_data.get("image_name", "mascot_think")
            has_ext = any(image_name.lower().endswith(ext) for ext in [".png", ".jpg", ".jpeg", ".gif"])
            if not has_ext:
                image_name += ".png"
            
            project_dir = os.path.dirname(os.path.abspath(output_path))
            custom_image_path = os.path.join(project_dir, "custom_images", image_name)
            if os.path.exists(custom_image_path):
                image_path = custom_image_path
            else:
                image_path = os.path.join(images_dir, image_name)
            
            if os.path.exists(image_path):
                slide.shapes.add_picture(image_path, Inches(6.8), Inches(1.4), width=Inches(5.4), height=Inches(4.8))

    # Create parent directories for output if they don't exist
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    
    prs.save(output_path)
    print(f"PowerPoint erfolgreich erstellt unter: {output_path}")

def main():
    parser = argparse.ArgumentParser(description="Automatische PPTX Generierung")
    parser.add_argument("--json-path", required=True, help="Pfad zur JSON Datei mit den Foliendaten")
    parser.add_argument("--output-path", required=True, help="Pfad für die zu erstellende PPTX-Datei")
    parser.add_argument("--images-dir", required=True, help="Verzeichnis, das die Standardbilder enthält")
    args = parser.parse_args()
    
    if not os.path.exists(args.json_path):
        print(f"Fehler: JSON-Datei {args.json_path} existiert nicht.")
        sys.exit(1)
        
    with open(args.json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
        
    create_presentation(data, args.output_path, args.images_dir)

if __name__ == "__main__":
    main()
