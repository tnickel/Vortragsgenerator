import sys
import os
import json
from pptx import Presentation

def import_pptx(pptx_path, output_json_path):
    if not os.path.exists(pptx_path):
        print(f"Error: PPTX file not found: {pptx_path}")
        sys.exit(1)
        
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Error opening PPTX file: {e}")
        sys.exit(1)
        
    slides_data = []
    
    for i, slide in enumerate(prs.slides):
        slide_number = i + 1
        title = ""
        bullets = []
        notes = ""
        
        # 1. Identify title shape (if any)
        title_shape = None
        # Try built-in title shape
        try:
            if slide.shapes.title:
                title_shape = slide.shapes.title
        except AttributeError:
            pass
            
        # If not found, look at placeholder shapes
        if not title_shape:
            for shape in slide.shapes:
                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    # Types: 1 is TITLE, 3 is CENTER_TITLE
                    if ph_type == 1 or ph_type == 3:
                        title_shape = shape
                        break
                        
        if title_shape and title_shape.has_text_frame:
            title = title_shape.text.strip()
            
        # 2. Extract speaker notes
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass
            
        # 3. Extract bullets/body text from other text shapes
        for shape in slide.shapes:
            if shape == title_shape:
                continue
            if shape.has_text_frame:
                tf = shape.text_frame
                for paragraph in tf.paragraphs:
                    text = paragraph.text.strip()
                    if text and text != title:
                        bullets.append(text)
                        
        slides_data.append({
            "slide_number": slide_number,
            "title": title or f"Folie {slide_number}",
            "bullets": bullets,
            "notes": notes,
            "image_name": "mascot_think.png"
        })
        
    project_data = {
        "title": os.path.basename(pptx_path).replace(".pptx", ""),
        "slides": slides_data
    }
    
    # Ensure parent directory of output_json_path exists
    output_dir = os.path.dirname(os.path.abspath(output_json_path))
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)
    
    with open(output_json_path, 'w', encoding='utf-8') as f:
        json.dump(project_data, f, ensure_ascii=False, indent=4)
        
    print(f"PPTX import successful. Metadata written to {output_json_path}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python import_pptx.py <pptx_path> <output_json_path>")
        sys.exit(1)
    import_pptx(sys.argv[1], sys.argv[2])
